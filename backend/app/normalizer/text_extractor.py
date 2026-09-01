"""
Text extraction orchestrator.

Routes to the right extraction path based on detected MIME type.
Supports: plain text, PDF, DOCX, XLSX, PPTX, HTML, images (via OCR).
All extraction is bounded by MAX_EXTRACTED_CHARS and OCR_TIMEOUT_SECONDS.
"""

import logging
from typing import Optional
import io
import os

logger = logging.getLogger(__name__)


class TextExtractor:
    """
    Orchestrates text extraction from various file formats.
    
    Routes to format-specific extractors and enforces size/time bounds.
    """
    
    def __init__(self, ocr_service, max_chars: int = 500000):
        """
        Initialize text extractor.
        
        Args:
            ocr_service: OCR service instance (real or fake)
            max_chars: Maximum extracted characters (truncate beyond this)
        """
        self.ocr_service = ocr_service
        self.max_chars = max_chars
    
    async def extract(
        self,
        content_bytes: bytes,
        mime_type: str,
        file_extension: Optional[str] = None,
        fixture_key: Optional[str] = None,
    ) -> str:
        """
        Extract text from file content.
        
        Routes to appropriate extractor based on MIME type.
        Truncates at max_chars to prevent unbounded memory use.
        
        Args:
            content_bytes: Raw file bytes
            mime_type: Detected MIME type
            file_extension: File extension (optional hint)
            fixture_key: Fixture identifier for tests (passed to fake OCR)
            
        Returns:
            Extracted text (bounded)
        """
        if not content_bytes:
            return ""
        
        # Normalize/resolve MIME type if missing or generic
        mime_type = (mime_type or "").strip().lower()
        if not mime_type or mime_type in ("application/octet-stream", "binary/octet-stream"):
            ext = (file_extension or "").lower().lstrip(".")
            if ext == "pdf":
                mime_type = "application/pdf"
            elif ext in ("docx", "doc"):
                mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            elif ext in ("xlsx", "xls"):
                mime_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            elif ext in ("pptx", "ppt"):
                mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            elif ext in ("txt", "text", "md", "csv", "json", "log"):
                mime_type = "text/plain"
            elif ext in ("html", "htm"):
                mime_type = "text/html"
            elif ext in ("png", "jpg", "jpeg", "webp", "tiff", "bmp"):
                mime_type = f"image/{ext}"

        # Route by MIME type
        mime_major = mime_type.split("/")[0] if "/" in mime_type else mime_type
        
        try:
            if mime_type == "text/html":
                # Must be checked BEFORE the generic text/* catch-all below
                text = self._extract_html(content_bytes)
            elif mime_type == "text/plain" or mime_major == "text":
                text = self._extract_plain_text(content_bytes)
            elif mime_type == "application/pdf":
                text = await self._extract_pdf(content_bytes, fixture_key)
            elif mime_type in (
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "application/msword",
            ):
                text = self._extract_docx(content_bytes)
            elif mime_type in (
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/vnd.ms-excel",
            ):
                text = self._extract_xlsx(content_bytes)
            elif mime_type in (
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "application/vnd.ms-powerpoint",
            ):
                text = self._extract_pptx(content_bytes)
            elif mime_major == "image":
                # Fall back to OCR for images
                text = self.ocr_service.extract_text(content_bytes, fixture_key=fixture_key)
            else:
                logger.warning(f"No text extractor for MIME type {mime_type}")
                text = ""
        except Exception as e:
            logger.error(f"Text extraction failed for {mime_type}: {e}")
            text = ""
        
        # Truncate to max_chars
        if len(text) > self.max_chars:
            logger.warning(
                f"Extracted text truncated from {len(text)} to {self.max_chars} chars"
            )
            text = text[: self.max_chars]
        
        return text.strip()
    
    def _extract_plain_text(self, content_bytes: bytes) -> str:
        """Extract text from plain text file."""
        try:
            return content_bytes.decode("utf-8")
        except UnicodeDecodeError:
            # Try latin-1 as fallback
            try:
                return content_bytes.decode("latin-1")
            except Exception as e:
                logger.error(f"Failed to decode plain text: {e}")
                return ""
    
    def _extract_html(self, content_bytes: bytes) -> str:
        """Extract text from HTML (strip tags)."""
        try:
            from bs4 import BeautifulSoup
            
            html = content_bytes.decode("utf-8", errors="ignore")
            soup = BeautifulSoup(html, "html.parser")
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            return soup.get_text(separator=" ", strip=True)
        except Exception as e:
            logger.error(f"HTML extraction failed: {e}")
            return ""
    
    async def _extract_pdf(self, content_bytes: bytes, fixture_key: Optional[str] = None) -> str:
        """
        Extract text from PDF using pdfplumber.
        
        Falls back to OCR for scanned/image-only pages, and pypdf if pdfplumber fails.
        """
        text_parts = []
        try:
            import pdfplumber
            
            with pdfplumber.open(io.BytesIO(content_bytes)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    
                    if page_text and page_text.strip():
                        text_parts.append(page_text)
                    else:
                        # Page has no text — likely scanned image
                        # Fall back to OCR (convert page to image first)
                        try:
                            page_image = page.to_image(resolution=150)
                            image_bytes = io.BytesIO()
                            page_image.original.save(image_bytes, format="PNG")
                            ocr_text = self.ocr_service.extract_text(
                                image_bytes.getvalue(), fixture_key=fixture_key
                            )
                            if ocr_text:
                                text_parts.append(ocr_text)
                        except Exception as ocr_error:
                            logger.warning(f"OCR fallback failed for PDF page: {ocr_error}")
        except Exception as e:
            logger.warning(f"PDF extraction with pdfplumber encountered error: {e}, trying pypdf fallback")
            try:
                import pypdf
                reader = pypdf.PdfReader(io.BytesIO(content_bytes))
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text_parts.append(page_text)
            except Exception as pypdf_err:
                logger.error(f"pypdf fallback also failed: {pypdf_err}")

        return "\n".join(text_parts)
    
    def _extract_docx(self, content_bytes: bytes) -> str:
        """Extract text from DOCX.

        Rule #3 — walks body paragraphs, tables, text boxes (w:txbxContent),
        headers, and footers.  Plain ``doc.paragraphs`` misses text in tables
        and text boxes, which is a confirmed gap for python-docx.
        """
        try:
            from docx import Document

            doc = Document(io.BytesIO(content_bytes))
            text_parts: list[str] = []

            # 1. Body paragraphs (reading order)
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    text_parts.append(text)

            # 2. Tables — iterate rows → cells → paragraphs
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for para in cell.paragraphs:
                            text = para.text.strip()
                            if text:
                                text_parts.append(text)

            # 3. Headers and footers — iterate sections
            for section in doc.sections:
                for header_footer in (
                    section.header,
                    section.footer,
                    getattr(section, "first_page_header", None),
                    getattr(section, "first_page_footer", None),
                    getattr(section, "even_page_header", None),
                    getattr(section, "even_page_footer", None),
                ):
                    if header_footer is None:
                        continue
                    try:
                        for para in header_footer.paragraphs:
                            text = para.text.strip()
                            if text:
                                text_parts.append(text)
                    except Exception:
                        pass

            # 4. Text boxes (w:txbxContent) — walk the XML directly
            try:
                from lxml import etree

                nsmap = {
                    "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
                    "mc": "http://schemas.openxmlformats.org/markup-compatibility/2006",
                    "wps": "http://schemas.microsoft.com/office/word/2010/wordprocessingShape",
                }
                body_xml = doc.element
                # Find all w:txbxContent elements (text boxes)
                for txbx in body_xml.iter(
                    "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}txbxContent"
                ):
                    for p_elem in txbx.iter(
                        "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"
                    ):
                        if p_elem.text and p_elem.text.strip():
                            text_parts.append(p_elem.text.strip())
            except ImportError:
                # lxml not available — skip text box extraction
                logger.debug("lxml not available, skipping text box extraction")
            except Exception as e:
                logger.debug(f"Text box extraction failed: {e}")

            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            return ""
    
    def _extract_xlsx(self, content_bytes: bytes) -> str:
        """Extract text from XLSX (cell values)."""
        try:
            from openpyxl import load_workbook
            
            wb = load_workbook(io.BytesIO(content_bytes), data_only=True)
            text_parts = []
            
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        text_parts.append(row_text)
            
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"XLSX extraction failed: {e}")
            return ""
    
    def _extract_pptx(self, content_bytes: bytes) -> str:
        """Extract text from PPTX (slide text)."""
        try:
            from pptx import Presentation
            
            prs = Presentation(io.BytesIO(content_bytes))
            text_parts = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)
            
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            return ""
